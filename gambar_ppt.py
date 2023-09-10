from pptx import Presentation
import os

def tambah_gambar(slide, path_gambar, kolom, baris):
    lebar_gambar = slide.slide_width // 10
    tinggi_gambar = slide.slide_height // 5
    x_pos = kolom * lebar_gambar
    y_pos = baris * tinggi_gambar
    slide.shapes.add_picture(path_gambar, x_pos, y_pos, width=lebar_gambar, height=tinggi_gambar)

def proses_folder(folder_utama, karakter):
    presentasi = Presentation()
    for sub_folder, _, files in os.walk(folder_utama):
        slide = presentasi.slides.add_slide(presentasi.slide_layouts[5])
        baris = 0
        kolom = 0
        for nama_file in files:
            if karakter in nama_file:
                path_gambar = os.path.join(sub_folder, nama_file)
                tambah_gambar(slide, path_gambar, kolom, baris)
                kolom += 1
                if kolom >= 10:
                    kolom = 0
                    baris += 1
                if baris >= 5:
                    baris = 0
                    slide = presentasi.slides.add_slide(presentasi.slide_layouts[5])
    presentasi.save('presentasi.pptx')

folder_utama = 'path_ke_folder_utama' 
karakter = 'A' # ganti dengan karakter yang mau difilter
proses_folder(folder_utama, karakter)